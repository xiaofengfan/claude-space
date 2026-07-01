# -*- coding: utf-8 -*-
"""Generate Claude Space product introduction PPT (dark tech theme)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------- Colors -----------------------------
BG_DARK      = RGBColor(0x0E, 0x0E, 0x14)
BG_PANEL     = RGBColor(0x16, 0x16, 0x20)
BG_CARD      = RGBColor(0x1C, 0x1C, 0x28)
BG_CARD_HI   = RGBColor(0x23, 0x23, 0x32)
LINE_COLOR   = RGBColor(0x2C, 0x2C, 0x3C)

ACC_PURPLE   = RGBColor(0x7C, 0x6C, 0xF0)
ACC_BLUE     = RGBColor(0x4E, 0x8D, 0xF3)
ACC_CYAN     = RGBColor(0x22, 0xD3, 0xEE)
ACC_GREEN    = RGBColor(0x34, 0xD9, 0x99)
ACC_ORANGE   = RGBColor(0xFF, 0xA3, 0x4D)
ACC_PINK     = RGBColor(0xF4, 0x6E, 0xB4)
ACC_YELLOW   = RGBColor(0xFD, 0xD6, 0x6B)
ACC_INDIGO   = RGBColor(0x82, 0x76, 0xE8)

TEXT_PRIMARY   = RGBColor(0xF1, 0xF1, 0xF6)
TEXT_SECONDARY = RGBColor(0xA6, 0xA6, 0xB8)
TEXT_MUTED     = RGBColor(0x6C, 0x6C, 0x7E)
TEXT_DIM       = RGBColor(0x50, 0x50, 0x60)

FONT_CN = 'Microsoft YaHei'
FONT_EN = 'Segoe UI'

# ----------------------------- Geometry ---------------------------
PAGE_W = Inches(13.333)
PAGE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = PAGE_W - 2 * MARGIN

# ----------------------------- Helpers ----------------------------
def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
             shape_type=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape_type, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w if line_w else Pt(0.75)
    return sp

def add_round(slide, x, y, w, h, fill=None, line=None, radius=0.08):
    sp = add_rect(slide, x, y, w, h, fill=fill, line=line, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp

def set_gradient(sp, c1, c2, angle=45):
    fill = sp.fill
    fill.gradient()
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = c2
    try:
        fill.gradient_angle = angle
    except Exception:
        pass

def add_text(slide, x, y, w, h, text, size=14, color=TEXT_PRIMARY, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_CN,
             spacing=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb

def add_runs(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    """runs: list of (text, dict-of-overrides)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    for text, ov in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(ov.get('size', 14))
        r.font.bold = ov.get('bold', False)
        r.font.italic = ov.get('italic', False)
        r.font.color.rgb = ov.get('color', TEXT_PRIMARY)
        r.font.name = ov.get('font', FONT_CN)
    return tb

def add_list(slide, x, y, w, h, items, size=13, color=TEXT_SECONDARY,
             bullet='\u25B8  ', bullet_color=None, spacing=1.25, gap=7):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    bc = bullet_color or ACC_CYAN
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        p.space_after = Pt(gap)
        rb = p.add_run()
        rb.text = bullet
        rb.font.size = Pt(size)
        rb.font.color.rgb = bc
        rb.font.bold = True
        rb.font.name = FONT_CN
        rt = p.add_run()
        rt.text = item
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
        rt.font.name = FONT_CN
    return tb

def add_header(slide, title, subtitle=None, page=None, total=15):
    # decorative vertical bar (gradient)
    bar = add_rect(slide, MARGIN, Inches(0.58), Inches(0.12), Inches(0.92),
                   fill=ACC_PURPLE)
    set_gradient(bar, ACC_PURPLE, ACC_CYAN, angle=90)
    add_text(slide, MARGIN + Inches(0.30), Inches(0.52),
             Inches(10.5), Inches(0.6), title, size=30, color=TEXT_PRIMARY, bold=True)
    if subtitle:
        add_text(slide, MARGIN + Inches(0.30), Inches(1.06),
                 Inches(11), Inches(0.4), subtitle, size=14, color=TEXT_SECONDARY)
    # brand top-right
    add_text(slide, Inches(10.6), Inches(0.6), Inches(2.13), Inches(0.4),
             'Claude Space', size=12, color=TEXT_MUTED, align=PP_ALIGN.RIGHT, font=FONT_EN)
    # footer
    fy = Inches(7.02)
    add_rect(slide, MARGIN, fy, CONTENT_W, Pt(0.75), fill=LINE_COLOR)
    add_text(slide, MARGIN, fy + Inches(0.06), Inches(6), Inches(0.3),
             'Claude Space  ·  AI 驱动的项目开发工作台', size=9, color=TEXT_DIM)
    if page:
        add_text(slide, Inches(11.0), fy + Inches(0.06), Inches(1.73), Inches(0.3),
                 f'{page:02d} / {total:02d}', size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT, font=FONT_EN)

def add_card(slide, x, y, w, h, title, desc=None, accent=ACC_PURPLE,
             title_size=15, desc_size=11.5, items=None, body_y_offset=0.62):
    add_round(slide, x, y, w, h, fill=BG_CARD, line=LINE_COLOR, radius=0.06)
    # left accent bar
    add_round(slide, x, y + Inches(0.14), Inches(0.07), h - Inches(0.28),
              fill=accent, radius=0.5)
    add_text(slide, x + Inches(0.28), y + Inches(0.20),
             w - Inches(0.45), Inches(0.4), title, size=title_size,
             color=TEXT_PRIMARY, bold=True)
    if desc:
        add_text(slide, x + Inches(0.28), y + Inches(body_y_offset),
                 w - Inches(0.45), h - Inches(body_y_offset + 0.2),
                 desc, size=desc_size, color=TEXT_SECONDARY, spacing=1.3)
    if items:
        add_list(slide, x + Inches(0.28), y + Inches(body_y_offset),
                 w - Inches(0.45), h - Inches(body_y_offset + 0.2),
                 items, size=desc_size, gap=4)


# ===================================================================
def build():
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    LOGO = r'E:\claudespace\claude-space\assets\icon.png'

    # ---------------- Slide 1: Cover ----------------
    s = add_slide(prs); set_bg(s, BG_DARK)
    # big gradient backdrop block
    bg = add_rect(s, 0, 0, PAGE_W, PAGE_H, fill=BG_PANEL)
    set_gradient(bg, RGBColor(0x14, 0x12, 0x28), RGBColor(0x08, 0x0C, 0x16), angle=60)
    # accent glow bars
    g1 = add_rect(s, Inches(-1), Inches(4.6), Inches(7), Inches(3.2), fill=BG_CARD)
    set_gradient(g1, ACC_PURPLE, BG_DARK, angle=35)
    g1.fill.transparency = 0  # no-op; keep
    g2 = add_rect(s, Inches(8.5), Inches(-1.2), Inches(6.5), Inches(4.2), fill=BG_CARD)
    set_gradient(g2, ACC_BLUE, BG_DARK, angle=120)
    # logo
    try:
        s.shapes.add_picture(LOGO, Inches(5.97), Inches(1.25), width=Inches(1.4))
    except Exception as e:
        print('logo skip', e)
    # title
    add_text(s, Inches(0.8), Inches(2.95), Inches(11.73), Inches(1.0),
             'Claude Space', size=58, color=TEXT_PRIMARY, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_EN)
    # slogan underline
    line = add_rect(s, Inches(5.67), Inches(4.05), Inches(2.0), Pt(3), fill=ACC_CYAN)
    set_gradient(line, ACC_PURPLE, ACC_CYAN, angle=0)
    add_text(s, Inches(0.8), Inches(4.25), Inches(11.73), Inches(0.6),
             'AI 驱动的项目开发工作台', size=26, color=TEXT_SECONDARY, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(4.95), Inches(11.73), Inches(0.5),
             '多模型  ·  多智能体  ·  多工作空间', size=18, color=ACC_CYAN,
             align=PP_ALIGN.CENTER, spacing=1.4)
    # version chip
    chip = add_round(s, Inches(5.92), Inches(5.85), Inches(1.5), Inches(0.5),
                     fill=BG_CARD, line=ACC_PURPLE, radius=0.5)
    add_text(s, Inches(5.92), Inches(5.9), Inches(1.5), Inches(0.42),
             'v1.1.7', size=15, color=TEXT_PRIMARY, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
    add_text(s, Inches(0.8), Inches(6.75), Inches(11.73), Inches(0.4),
             '将 Claude Code CLI 升级为图形化多智能体协作平台', size=12,
             color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ---------------- Slide 2: Positioning ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '产品定位', '从一个聊天窗口，进化为 AI 指挥中心', page=2)
    # evolution arrow: 3 stages
    stages = [
        ('AI Chat', '问答式助手', '单轮对话 · 文本输出', ACC_BLUE),
        ('Claude Code CLI', '终端 AI', '命令行 · 单智能体', ACC_PURPLE),
        ('Claude Space', 'AI 工作站', '图形化 · 多智能体协作', ACC_CYAN),
    ]
    sx = Inches(0.6); sy = Inches(2.0); sw = Inches(3.6); sh = Inches(1.9); gap = Inches(0.35)
    for i, (t, sub, d, c) in enumerate(stages):
        x = sx + i * (sw + gap)
        add_round(s, x, sy, sw, sh, fill=BG_CARD, line=LINE_COLOR, radius=0.08)
        add_rect(s, x, sy, sw, Inches(0.08), fill=c)
        add_text(s, x, sy + Inches(0.35), sw, Inches(0.5), t, size=20, bold=True,
                 color=TEXT_PRIMARY, align=PP_ALIGN.CENTER, font=FONT_EN)
        add_text(s, x, sy + Inches(0.95), sw, Inches(0.4), sub, size=14, bold=True,
                 color=c, align=PP_ALIGN.CENTER)
        add_text(s, x, sy + Inches(1.4), sw, Inches(0.4), d, size=11,
                 color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        if i < 2:
            ax = x + sw + Inches(0.02)
            add_text(s, ax, sy + Inches(0.7), gap, Inches(0.5), '\u2192', size=24,
                     color=TEXT_MUTED, align=PP_ALIGN.CENTER, bold=True)
    # bottom contrast cards
    by = Inches(4.35); bh = Inches(2.2)
    # left: traditional
    add_round(s, Inches(0.6), by, Inches(5.86), bh, fill=BG_PANEL, line=LINE_COLOR, radius=0.05)
    add_text(s, Inches(0.9), by + Inches(0.22), Inches(5), Inches(0.4),
             '传统方式局限', size=15, bold=True, color=ACC_PINK)
    add_list(s, Inches(0.95), by + Inches(0.72), Inches(5.2), Inches(1.4),
             ['单轮问答，难以追踪长任务', '纯文本输出，工具调用不可视',
              '单智能体，无法团队协作', '缺乏审批，敏感操作不可控'],
             size=12.5, color=TEXT_SECONDARY, bullet='\u2717  ', bullet_color=ACC_PINK, gap=5)
    # right: claude space
    add_round(s, Inches(6.87), by, Inches(5.86), bh, fill=BG_CARD_HI, line=ACC_CYAN, radius=0.05)
    add_text(s, Inches(7.17), by + Inches(0.22), Inches(5), Inches(0.4),
             'Claude Space 优势', size=15, bold=True, color=ACC_CYAN)
    add_list(s, Inches(7.22), by + Inches(0.72), Inches(5.2), Inches(1.4),
             ['多模型接入，多 Tab 会话并行', '图形化任务看板，实时监控',
              '9人虚拟团队，@mention 协作', '敏感操作审批，全流程可追溯'],
             size=12.5, color=TEXT_PRIMARY, bullet='\u2713  ', bullet_color=ACC_GREEN, gap=5)

    # ---------------- Slide 3: Capabilities overview ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '核心能力概览', '八大模块，构建完整的 AI 开发工作流', page=3)
    caps = [
        ('多模型 AI 对话', 'Claude / DeepSeek / OpenAI / 自定义，流式渲染与多 Tab 会话', ACC_PURPLE),
        ('多智能体协作', '9人虚拟团队，5种角色原型，队列并发调度', ACC_BLUE),
        ('项目管理', '自动扫描、多工作空间、递归文件树与代码编辑', ACC_CYAN),
        ('集成终端', 'node-pty + xterm.js 本地终端，WebGL 加速渲染', ACC_GREEN),
        ('任务看板', '三列看板、AI 自动提取任务、实时监控', ACC_ORANGE),
        ('Git 集成', '状态/差异/提交/分支，可视化的版本控制', ACC_PINK),
        ('SSH 远程开发', '远程终端、SFTP 浏览、一键部署上线', ACC_INDIGO),
        ('记忆与知识', '项目隔离记忆、知识库 CRUD、规则模板', ACC_YELLOW),
    ]
    cols = 4; rows = 2
    gx = Inches(0.6); gy = Inches(1.85)
    cw = Inches(2.9); ch = Inches(2.35)
    gxgap = Inches(0.16); gygap = Inches(0.22)
    for i, (t, d, c) in enumerate(caps):
        r = i // cols; col = i % cols
        x = gx + col * (cw + gxgap)
        y = gy + r * (ch + gygap)
        add_card(s, x, y, cw, ch, t, desc=d, accent=c,
                 title_size=15, desc_size=11, body_y_offset=0.66)

    # ---------------- Slide 4: Tech architecture ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '技术架构', 'Electron 主进程 + React 渲染进程，安全隔离通信', page=4)
    # main process panel
    mx = Inches(0.6); my = Inches(1.95); mw = Inches(4.0); mh = Inches(4.4)
    add_round(s, mx, my, mw, mh, fill=BG_PANEL, line=LINE_COLOR, radius=0.04)
    add_rect(s, mx, my, mw, Inches(0.5), fill=ACC_PURPLE)
    add_text(s, mx, my, mw, Inches(0.5), '主进程 · Main (Node.js)', size=14, bold=True,
             color=TEXT_PRIMARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    main_mods = [
        ('claudeProcess', 'Claude CLI 子进程与 JSONL 流解析'),
        ('agentPool', '多智能体池，队列并发（上限3）'),
        ('terminalProcess', '本地 PTY 终端与会话同步'),
        ('sshService', 'SSH 连接池、SFTP、部署管线'),
        ('connectionService', 'CLI 检测、API 探测、健康检查'),
        ('main.ts', '170+ IPC 处理器，任务/设置持久化'),
    ]
    yy = my + Inches(0.7)
    for name, d in main_mods:
        add_text(s, mx + Inches(0.3), yy, Inches(3.5), Inches(0.3), name,
                 size=12.5, bold=True, color=ACC_CYAN, font=FONT_EN)
        add_text(s, mx + Inches(0.3), yy + Inches(0.3), Inches(3.5), Inches(0.35), d,
                 size=10.5, color=TEXT_SECONDARY, spacing=1.2)
        yy = yy + Inches(0.6)
    # bridge in the middle
    bx = Inches(4.8); bw = Inches(3.73)
    add_round(s, bx, Inches(3.55), bw, Inches(1.2), fill=BG_CARD_HI, line=ACC_BLUE, radius=0.1)
    add_text(s, bx, Inches(3.62), bw, Inches(0.4), 'contextBridge / IPC',
             size=15, bold=True, color=ACC_BLUE, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(s, bx, Inches(4.05), bw, Inches(0.4), '160+ 安全方法 · 零信任隔离',
             size=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
    add_text(s, bx, Inches(4.4), bw, Inches(0.35), 'credentials 仅驻留主进程',
             size=10.5, color=ACC_GREEN, align=PP_ALIGN.CENTER)
    # arrows
    add_text(s, Inches(4.6), Inches(3.95), Inches(0.25), Inches(0.4), '\u2192',
             size=22, color=TEXT_MUTED, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(8.5), Inches(3.95), Inches(0.25), Inches(0.4), '\u2190',
             size=22, color=TEXT_MUTED, align=PP_ALIGN.CENTER, bold=True)
    # renderer panel
    rx = Inches(8.73); rw = Inches(4.0)
    add_round(s, rx, my, rw, mh, fill=BG_PANEL, line=LINE_COLOR, radius=0.04)
    add_rect(s, rx, my, rw, Inches(0.5), fill=ACC_CYAN)
    add_text(s, rx, my, rw, Inches(0.5), '渲染进程 · Renderer (React)', size=14,
             bold=True, color=BG_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    panels = [
        ('左侧栏', '文件树 / 会话 / Git / 记忆 / 规则'),
        ('中央区', '聊天 / 终端 / 欢迎页'),
        ('右侧栏', '看板 / 监控 / Pixel Office / 部署'),
        ('底部', '状态栏（模型 · Token · 成本）'),
    ]
    yy = my + Inches(0.75)
    for name, d in panels:
        add_text(s, rx + Inches(0.3), yy, Inches(3.5), Inches(0.3), name,
                 size=12.5, bold=True, color=ACC_PURPLE)
        add_text(s, rx + Inches(0.3), yy + Inches(0.3), Inches(3.5), Inches(0.35), d,
                 size=10.5, color=TEXT_SECONDARY, spacing=1.2)
        yy = yy + Inches(0.85)
    # tech stack chips bottom
    sy = Inches(6.5)
    chips = ['Electron 28', 'React 18', 'TypeScript 5', 'Vite 5', 'xterm.js 5', 'node-pty', 'ssh2', 'electron-builder 24']
    cx = Inches(0.6)
    for ch in chips:
        w = Inches(0.16 * len(ch) + 0.5)
        add_round(s, cx, sy, w, Inches(0.36), fill=BG_CARD, line=LINE_COLOR, radius=0.5)
        add_text(s, cx, sy, w, Inches(0.36), ch, size=10, color=TEXT_SECONDARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
        cx = cx + w + Inches(0.12)

    # ---------------- Slide 5: Multi-model chat ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '多模型 AI 对话', '统一接入多家模型，流畅的流式交互体验', page=5)
    # left feature list
    add_round(s, Inches(0.6), Inches(1.95), Inches(6.2), Inches(4.4),
              fill=BG_PANEL, line=LINE_COLOR, radius=0.04)
    add_text(s, Inches(0.9), Inches(2.2), Inches(5.5), Inches(0.4),
             '对话能力', size=16, bold=True, color=ACC_PURPLE)
    add_list(s, Inches(0.95), Inches(2.75), Inches(5.6), Inches(3.4),
             ['支持 Claude / DeepSeek / OpenAI / 自定义模型，一键切换',
              '流式消息渲染，实时展示 AI 思考过程（Thinking Block）',
              '工具调用卡片：Bash / Read / Write / Edit / Grep / Glob',
              '多 Tab 会话管理，每个会话独立历史与上下文',
              '图片粘贴 / 拖拽支持，多模态输入',
              '@mention 指定智能体响应，群聊协作'],
             size=13, color=TEXT_SECONDARY, bullet='\u25B8  ', bullet_color=ACC_CYAN, gap=10)
    # right: model cards
    models = [
        ('Claude', 'Anthropic 官方 · Claude Code CLI', ACC_PURPLE),
        ('DeepSeek', '国产高性价比 · 自定义端点', ACC_BLUE),
        ('OpenAI', 'GPT 系列 · 兼容接口', ACC_GREEN),
        ('Custom', '任意 OpenAI 兼容服务', ACC_ORANGE),
    ]
    rx = Inches(7.1); rw = Inches(5.63)
    add_text(s, rx, Inches(2.0), rw, Inches(0.4), '已支持模型', size=16, bold=True, color=ACC_CYAN)
    my2 = Inches(2.55); mcw = Inches(2.72); mch = Inches(1.15)
    for i, (n, d, c) in enumerate(models):
        col = i % 2; row = i // 2
        x = rx + col * (mcw + Inches(0.19))
        y = my2 + row * (mch + Inches(0.2))
        add_round(s, x, y, mcw, mch, fill=BG_CARD, line=LINE_COLOR, radius=0.08)
        add_rect(s, x, y + Inches(0.16), Inches(0.07), mch - Inches(0.32), fill=c, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x + Inches(0.28), y + Inches(0.22), mcw - Inches(0.4), Inches(0.4),
                 n, size=16, bold=True, color=TEXT_PRIMARY, font=FONT_EN)
        add_text(s, x + Inches(0.28), y + Inches(0.66), mcw - Inches(0.4), Inches(0.4),
                 d, size=10.5, color=TEXT_SECONDARY, spacing=1.2)
    # chat mock strip
    cy = Inches(5.15)
    add_round(s, rx, cy, rw, Inches(1.2), fill=BG_CARD_HI, line=LINE_COLOR, radius=0.06)
    add_text(s, rx + Inches(0.25), cy + Inches(0.12), rw, Inches(0.35),
             '用户', size=11, bold=True, color=ACC_CYAN)
    add_text(s, rx + Inches(0.25), cy + Inches(0.42), rw - Inches(0.5), Inches(0.3),
             '帮我看下这个文件的结构', size=11.5, color=TEXT_PRIMARY)
    add_text(s, rx + Inches(0.25), cy + Inches(0.78), rw, Inches(0.3),
             'AI · 正在调用 Read 工具...', size=10.5, color=TEXT_MUTED, italic=True)

    # ---------------- Slide 6: Multi-agent ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '多智能体团队协作', '5种角色原型，9人虚拟团队，并行作业', page=6)
    # left: 5 archetypes
    add_text(s, Inches(0.6), Inches(1.85), Inches(5.8), Inches(0.4),
             '5 种智能体原型', size=15, bold=True, color=ACC_PURPLE)
    archetypes = [
        ('Coordinator', '协调员', '项目/产品管理 · 任务分配', ACC_PURPLE),
        ('Architect', '架构师', '系统设计 · 技术选型', ACC_BLUE),
        ('Implementer', '实现者', '编码实现 · 功能开发', ACC_GREEN),
        ('SecurityReviewer', '安全审查员', '测试 · 安全审计', ACC_ORANGE),
        ('CodeExplorer', '代码探索者', '代码分析 · 探索', ACC_CYAN),
    ]
    ay = Inches(2.35); ah = Inches(0.78)
    for i, (en, cn, d, c) in enumerate(archetypes):
        y = ay + i * (ah + Inches(0.08))
        add_round(s, Inches(0.6), y, Inches(5.8), ah, fill=BG_CARD, line=LINE_COLOR, radius=0.1)
        add_rect(s, Inches(0.6), y + Inches(0.16), Inches(0.07), ah - Inches(0.32), fill=c)
        add_text(s, Inches(0.85), y + Inches(0.1), Inches(2.3), Inches(0.3), en,
                 size=13, bold=True, color=c, font=FONT_EN)
        add_text(s, Inches(0.85), y + Inches(0.38), Inches(2.3), Inches(0.3), cn,
                 size=11.5, bold=True, color=TEXT_PRIMARY)
        add_text(s, Inches(3.1), y + Inches(0.24), Inches(3.2), Inches(0.4), d,
                 size=11, color=TEXT_SECONDARY, spacing=1.2)
    # right: 9 member team grid
    rx = Inches(6.75); rw = Inches(5.98)
    add_text(s, rx, Inches(1.85), rw, Inches(0.4),
             '9 人默认团队', size=15, bold=True, color=ACC_CYAN)
    members = [
        ('王PM', '协调员', ACC_PURPLE), ('李PM', '协调员', ACC_PURPLE),
        ('张架构', '架构师', ACC_BLUE), ('赵高工', '实现者', ACC_GREEN),
        ('钱开发', '实现者', ACC_GREEN), ('孙开发', '实现者', ACC_GREEN),
        ('周QA', '审查员', ACC_ORANGE), ('吴审查', '审查员', ACC_ORANGE),
        ('Claude', 'AI', ACC_CYAN),
    ]
    mcw = Inches(1.86); mch = Inches(1.05)
    mx0 = rx; my0 = Inches(2.4)
    for i, (n, role, c) in enumerate(members):
        col = i % 3; row = i // 3
        x = mx0 + col * (mcw + Inches(0.12))
        y = my0 + row * (mch + Inches(0.14))
        add_round(s, x, y, mcw, mch, fill=BG_CARD, line=LINE_COLOR, radius=0.1)
        # avatar dot
        dot = add_rect(s, x + Inches(0.15), y + Inches(0.18), Inches(0.34), Inches(0.34),
                       fill=c, shape_type=MSO_SHAPE.OVAL)
        add_text(s, x + Inches(0.58), y + Inches(0.15), mcw - Inches(0.7), Inches(0.35),
                 n, size=12, bold=True, color=TEXT_PRIMARY)
        add_text(s, x + Inches(0.58), y + Inches(0.45), mcw - Inches(0.7), Inches(0.3),
                 role, size=9.5, color=c)
        # status
        add_rect(s, x + Inches(0.15), y + Inches(0.68), Inches(0.1), Inches(0.1),
                 fill=ACC_GREEN, shape_type=MSO_SHAPE.OVAL)
        add_text(s, x + Inches(0.32), y + Inches(0.62), Inches(1), Inches(0.25),
                 '就绪', size=9, color=TEXT_MUTED)
    # note bottom
    add_round(s, rx, Inches(6.05), rw, Inches(0.55), fill=BG_PANEL, line=ACC_PURPLE, radius=0.15)
    add_text(s, rx + Inches(0.2), Inches(6.08), rw, Inches(0.5),
             '队列调度 · 最大 3 并发 · @mention 指派 · 上下文传递',
             size=11.5, color=TEXT_SECONDARY, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)

    # ---------------- Slide 7: Pixel Office ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, 'Pixel Office 虚拟办公室', '等距像素风工作空间，团队状态可视化', page=7)
    # feature cards top row
    feats = [
        ('等距像素风', 'Isometric 像素艺术风格虚拟办公场景', ACC_PURPLE),
        ('9 位成员', '团队成员在工位实时可视化呈现', ACC_BLUE),
        ('状态动画', '工作 / 忙碌 / 空闲 动画与进度条', ACC_GREEN),
        ('思考气泡', '随机气泡展示成员当前关注点', ACC_CYAN),
    ]
    fx = Inches(0.6); fy = Inches(1.95); fw = Inches(2.9); fh = Inches(1.7); fgap = Inches(0.16)
    for i, (t, d, c) in enumerate(feats):
        x = fx + i * (fw + fgap)
        add_card(s, x, fy, fw, fh, t, desc=d, accent=c, title_size=14, desc_size=11, body_y_offset=0.62)
    # bottom: simulated office grid
    oy = Inches(4.05); ow = Inches(12.13); oh = Inches(2.55)
    add_round(s, Inches(0.6), oy, ow, oh, fill=BG_PANEL, line=LINE_COLOR, radius=0.03)
    add_text(s, Inches(0.9), oy + Inches(0.2), Inches(6), Inches(0.4),
             '虚拟办公场景示意', size=14, bold=True, color=ACC_CYAN)
    # desks grid 3x3
    descs = [
        ('王PM', '编写需求', ACC_PURPLE, 'busy'),
        ('李PM', '评审中', ACC_PURPLE, 'work'),
        ('张架构', '设计模块', ACC_BLUE, 'busy'),
        ('赵高工', '编码中', ACC_GREEN, 'busy'),
        ('钱开发', '空闲', ACC_GREEN, 'idle'),
        ('孙开发', '编码中', ACC_GREEN, 'work'),
        ('周QA', '测试中', ACC_ORANGE, 'busy'),
        ('吴审查', '审查中', ACC_ORANGE, 'work'),
        ('Claude', '思考中', ACC_CYAN, 'busy'),
    ]
    status_color = {'busy': ACC_PINK, 'work': ACC_GREEN, 'idle': TEXT_MUTED}
    status_text = {'busy': '忙碌', 'work': '工作中', 'idle': '空闲'}
    dx0 = Inches(1.0); dy0 = oy + Inches(0.7); dw = Inches(3.75); dh = Inches(0.55); dgapx = Inches(0.12); dgapy = Inches(0.12)
    for i, (n, task, c, st) in enumerate(descs):
        col = i % 3; row = i // 3
        x = dx0 + col * (dw + dgapx)
        y = dy0 + row * (dh + dgapy)
        add_round(s, x, y, dw, dh, fill=BG_CARD, line=LINE_COLOR, radius=0.15)
        add_rect(s, x + Inches(0.12), y + Inches(0.15), Inches(0.26), Inches(0.26),
                 fill=c, shape_type=MSO_SHAPE.OVAL)
        add_text(s, x + Inches(0.48), y + Inches(0.07), Inches(1.2), Inches(0.25),
                 n, size=11, bold=True, color=TEXT_PRIMARY)
        add_text(s, x + Inches(0.48), y + Inches(0.3), Inches(1.6), Inches(0.22),
                 task, size=9.5, color=TEXT_SECONDARY)
        sc = status_color[st]
        add_rect(s, x + Inches(2.7), y + Inches(0.2), Inches(0.12), Inches(0.12),
                 fill=sc, shape_type=MSO_SHAPE.OVAL)
        add_text(s, x + Inches(2.88), y + Inches(0.1), Inches(0.8), Inches(0.3),
                 status_text[st], size=9.5, color=sc, bold=True)

    # ---------------- Slide 8: Project & file mgmt ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '项目管理与文件操作', '多工作空间、文件树、多 Tab 代码编辑', page=8)
    items = [
        ('多工作空间管理', ['独立工作目录与配置', '工作空间隔离的进程池', '一键切换，状态保留'], ACC_PURPLE),
        ('项目扫描与切换', ['自动扫描 Claude 项目', '最近项目快速打开', '欢迎页引导选择'], ACC_BLUE),
        ('递归文件树', ['目录递归展开', 'Git 状态图标', '快速定位文件'], ACC_CYAN),
        ('代码编辑器', ['多 Tab 标签页', '50+ 语言语法高亮', '二进制文件识别'], ACC_GREEN),
    ]
    cols = 4; cw = Inches(2.9); ch = Inches(4.4); gap = Inches(0.16); gx = Inches(0.6); gy = Inches(1.95)
    for i, (t, lst, c) in enumerate(items):
        x = gx + i * (cw + gap)
        add_card(s, x, gy, cw, ch, t, items=lst, accent=c, title_size=15,
                 desc_size=11.5, body_y_offset=0.7)

    # ---------------- Slide 9: Terminal & SSH ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '集成终端与 SSH 远程', '本地 PTY 终端，远程 SSH 开发与一键部署', page=9)
    # left: local terminal
    lx = Inches(0.6); ly = Inches(1.95); lw = Inches(5.86); lh = Inches(4.4)
    add_round(s, lx, ly, lw, lh, fill=BG_PANEL, line=ACC_GREEN, radius=0.04)
    add_rect(s, lx, ly, lw, Inches(0.5), fill=ACC_GREEN)
    add_text(s, lx, ly, lw, Inches(0.5), '本地终端', size=14, bold=True,
             color=BG_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_list(s, lx + Inches(0.35), ly + Inches(0.75), lw - Inches(0.6), Inches(3.4),
             ['node-pty 驱动的真实 PTY 终端',
              'xterm.js 5.x + WebGL 加速渲染',
              '暗色/亮色双主题适配',
              'Claude 会话 JSONL 自动同步',
              'Web 链接可点击，自动适配尺寸',
              'Shell 降级回退，保障可用性'],
             size=12.5, color=TEXT_SECONDARY, bullet='\u25B8  ', bullet_color=ACC_GREEN, gap=11)
    # right: SSH remote
    rx = Inches(6.87); rw = Inches(5.86)
    add_round(s, rx, ly, rw, lh, fill=BG_PANEL, line=ACC_BLUE, radius=0.04)
    add_rect(s, rx, ly, rw, Inches(0.5), fill=ACC_BLUE)
    add_text(s, rx, ly, rw, Inches(0.5), 'SSH 远程开发', size=14, bold=True,
             color=TEXT_PRIMARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ssh_items = [
        ('远程终端', 'ssh2 shell，ITerminalProvider 抽象'),
        ('远程文件浏览', 'SFTP 递归目录树浏览'),
        ('一键部署', 'SFTP 上传 + 预/后部署命令'),
        ('连接管理', '密码 / 密钥认证，连接池复用'),
        ('安全验证', 'TOFU 指纹首次信任校验'),
        ('远程 Git', '远程仓库操作面板'),
    ]
    yy = ly + Inches(0.72)
    for t, d in ssh_items:
        add_rect(s, rx + Inches(0.3), yy + Inches(0.08), Inches(0.1), Inches(0.1),
                 fill=ACC_CYAN, shape_type=MSO_SHAPE.OVAL)
        add_text(s, rx + Inches(0.55), yy, Inches(2.1), Inches(0.3), t,
                 size=12.5, bold=True, color=TEXT_PRIMARY)
        add_text(s, rx + Inches(2.55), yy, rw - Inches(2.8), Inches(0.3), d,
                 size=11, color=TEXT_SECONDARY)
        yy = yy + Inches(0.58)

    # ---------------- Slide 10: Task & approval ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '任务管理与审批', '看板追踪、实时监控、敏感操作审批', page=10)
    # kanban mock
    ky = Inches(1.95); kh = Inches(2.5)
    cols3 = [
        ('待办', TEXT_MUTED, ['分析需求结构', '设计数据模型'], TEXT_MUTED),
        ('进行中', ACC_ORANGE, ['编写核心模块', '集成终端调试'], ACC_ORANGE),
        ('已完成', ACC_GREEN, ['项目扫描', '多模型接入'], ACC_GREEN),
    ]
    kcw = Inches(3.93); kgap = Inches(0.17); kx = Inches(0.6)
    for i, (name, _, tasks, c) in enumerate(cols3):
        x = kx + i * (kcw + kgap)
        add_round(s, x, ky, kcw, kh, fill=BG_PANEL, line=LINE_COLOR, radius=0.04)
        add_rect(s, x, ky, Inches(0.07), kh, fill=c)
        add_text(s, x + Inches(0.25), ky + Inches(0.18), kcw, Inches(0.35),
                 name, size=14, bold=True, color=c)
        add_text(s, x + kcw - Inches(0.6), ky + Inches(0.18), Inches(0.5), Inches(0.35),
                 str(len(tasks)), size=14, bold=True, color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, font=FONT_EN)
        ty = ky + Inches(0.62)
        for t in tasks:
            add_round(s, x + Inches(0.2), ty, kcw - Inches(0.4), Inches(0.62),
                      fill=BG_CARD, line=LINE_COLOR, radius=0.12)
            add_rect(s, x + Inches(0.32), ty + Inches(0.22), Inches(0.18), Inches(0.18),
                     fill=c, shape_type=MSO_SHAPE.OVAL)
            add_text(s, x + Inches(0.6), ty + Inches(0.16), kcw - Inches(0.8), Inches(0.35),
                     t, size=11.5, color=TEXT_PRIMARY)
            ty = ty + Inches(0.72)
    # bottom: approval flow
    ay = Inches(4.7); aw = Inches(12.13); ah = Inches(1.85)
    add_round(s, Inches(0.6), ay, aw, ah, fill=BG_CARD, line=ACC_PINK, radius=0.03)
    add_text(s, Inches(0.9), ay + Inches(0.18), Inches(5), Inches(0.4),
             '敏感操作审批流程', size=15, bold=True, color=ACC_PINK)
    flow = ['AI 调用工具', '弹出审批', '用户决策', '执行/拒绝', '日志持久化']
    fcw = Inches(2.05); fx = Inches(0.95); fyy = ay + Inches(0.75)
    for i, step in enumerate(flow):
        x = fx + i * Inches(2.32)
        add_round(s, x, fyy, fcw, Inches(0.7), fill=BG_PANEL, line=LINE_COLOR, radius=0.15)
        add_text(s, x, fyy, fcw, Inches(0.7), step, size=11.5, bold=True,
                 color=TEXT_PRIMARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(flow) - 1:
            add_text(s, x + fcw, fyy, Inches(0.27), Inches(0.7), '\u2192', size=16,
                     color=TEXT_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)

    # ---------------- Slide 11: Memory & knowledge ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '记忆与知识管理', '项目隔离记忆，沉淀团队知识资产', page=11)
    items = [
        ('项目记忆系统', ['按项目隔离存储', '自动记忆会话与提交', '详情查看与编辑'], ACC_PURPLE),
        ('知识库管理', ['完整 CRUD 操作', '分类组织知识条目', '会话提取为知识'], ACC_BLUE),
        ('规则管理', ['CLAUDE.md 规则维护', '项目约定可视化', '锁定关键规则'], ACC_CYAN),
        ('规则模板库', ['7 大分类模板', '代码审查 / 安全', '约定 / API / 测试'], ACC_GREEN),
    ]
    cols = 4; cw = Inches(2.9); ch = Inches(4.4); gap = Inches(0.16); gx = Inches(0.6); gy = Inches(1.95)
    for i, (t, lst, c) in enumerate(items):
        x = gx + i * (cw + gap)
        add_card(s, x, gy, cw, ch, t, items=lst, accent=c, title_size=15,
                 desc_size=11.5, body_y_offset=0.7)

    # ---------------- Slide 12: Security ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '安全架构', '零信任设计，凭证不离开主进程', page=12)
    secs = [
        ('零信任渲染器', '渲染进程视为不可信，所有特权操作经主进程校验', ACC_PURPLE),
        ('contextBridge 隔离', '160+ 方法经安全桥暴露，禁用 Node 直接访问', ACC_BLUE),
        ('凭证驻留主进程', 'API Key / SSH 密钥仅存主进程，脱敏后传递', ACC_CYAN),
        ('TOFU 指纹验证', 'SSH 首次信任校验，防止中间人攻击', ACC_GREEN),
        ('部署排除机制', 'DEFAULT_DEPLOY_EXCLUDES 过滤敏感文件', ACC_ORANGE),
        ('敏感操作审批', 'Bash/Write/Edit 等需用户确认，全流程留痕', ACC_PINK),
    ]
    cols = 3; rows = 2
    gx = Inches(0.6); gy = Inches(1.95)
    cw = Inches(3.93); ch = Inches(2.1); gxgap = Inches(0.17); gygap = Inches(0.2)
    for i, (t, d, c) in enumerate(secs):
        r = i // cols; col = i % cols
        x = gx + col * (cw + gxgap)
        y = gy + r * (ch + gygap)
        add_card(s, x, y, cw, ch, t, desc=d, accent=c, title_size=15,
                 desc_size=12, body_y_offset=0.7)
    # bottom banner
    by = Inches(6.45)
    add_round(s, Inches(0.6), by, Inches(12.13), Inches(0.5), fill=BG_CARD_HI, line=ACC_GREEN, radius=0.15)
    add_text(s, Inches(0.6), by, Inches(12.13), Inches(0.5),
             '\u2713  安全原则：最小权限 · 主进程兜底 · 凭证脱敏 · 操作可审计',
             size=12.5, bold=True, color=ACC_GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---------------- Slide 13: Roadmap ----------------
    s = add_slide(prs); set_bg(s)
    add_header(s, '版本路线图', '从单机工作站，迈向数字公司操作系统', page=13)
    phases = [
        ('Phase 1', 'v1.0 – 1.1', '基础工作站', '已完成全部核心功能', True, ACC_GREEN),
        ('Phase 2', 'v1.2', 'Claw 自主引擎', '事件驱动开发循环 · 需求解析 · CI/CD', False, ACC_PURPLE),
        ('Phase 3', 'v1.3', '云端多方协作', 'GitHub/GitLab 连接 · 边缘-云端协作', False, ACC_BLUE),
        ('Phase 4', 'v1.4', '多用户协作', '账号 · 实时同步 · 团队消息 · 权限', False, ACC_CYAN),
        ('Phase 5', 'v1.5', '跨客户端消息', 'WebSocket/WebRTC · Web 客户端', False, ACC_ORANGE),
    ]
    py = Inches(1.95); pw = Inches(12.13); ph = Inches(0.88); pgap = Inches(0.06)
    for i, (p, v, name, d, done, c) in enumerate(phases):
        y = py + i * (ph + pgap)
        add_round(s, Inches(0.6), y, pw, ph, fill=BG_CARD_HI if done else BG_PANEL,
                  line=c if done else LINE_COLOR, radius=0.06)
        # status dot
        add_rect(s, Inches(0.85), y + Inches(0.34), Inches(0.2), Inches(0.2),
                 fill=c, shape_type=MSO_SHAPE.OVAL)
        mark = '\u2713' if done else '\u25CB'
        add_text(s, Inches(0.82), y + Inches(0.28), Inches(0.3), Inches(0.3), mark,
                 size=13, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.3), y + Inches(0.12), Inches(1.1), Inches(0.3), p,
                 size=12, bold=True, color=c, font=FONT_EN)
        add_text(s, Inches(1.3), y + Inches(0.42), Inches(1.1), Inches(0.3), v,
                 size=10.5, color=TEXT_MUTED, font=FONT_EN)
        add_text(s, Inches(2.6), y + Inches(0.16), Inches(2.6), Inches(0.35), name,
                 size=14, bold=True, color=TEXT_PRIMARY)
        add_text(s, Inches(2.6), y + Inches(0.48), Inches(7), Inches(0.3), d,
                 size=11, color=TEXT_SECONDARY)
        if done:
            add_round(s, Inches(11.0), y + Inches(0.28), Inches(1.5), Inches(0.34),
                      fill=c, radius=0.5)
            add_text(s, Inches(11.0), y + Inches(0.28), Inches(1.5), Inches(0.34),
                     '已完成', size=10.5, bold=True, color=BG_DARK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---------------- Slide 14: Vision ----------------
    s = add_slide(prs); set_bg(s, BG_PANEL)
    # gradient backdrop
    bg = add_rect(s, 0, 0, PAGE_W, PAGE_H, fill=BG_PANEL)
    set_gradient(bg, RGBColor(0x14, 0x10, 0x26), RGBColor(0x08, 0x0C, 0x16), angle=60)
    add_text(s, Inches(0.8), Inches(0.7), Inches(11.73), Inches(0.5),
             '终极愿景', size=16, bold=True, color=ACC_CYAN, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(1.3), Inches(11.73), Inches(1.0),
             '个人 / 团队 数字公司 OS', size=44, bold=True, color=TEXT_PRIMARY,
             align=PP_ALIGN.CENTER)
    line = add_rect(s, Inches(5.67), Inches(2.45), Inches(2.0), Pt(3), fill=ACC_PURPLE)
    set_gradient(line, ACC_PURPLE, ACC_CYAN, angle=0)
    add_text(s, Inches(0.8), Inches(2.7), Inches(11.73), Inches(0.5),
             '24×7 AI 数字员工 · 自进化多模型开发中心', size=18, color=TEXT_SECONDARY,
             align=PP_ALIGN.CENTER)
    # 4 vision cards
    visions = [
        ('多智能体平台', 'Claude Space + OpenClaw + Hermes 协同', ACC_PURPLE),
        ('共享记忆层', '向量数据库 + 知识图谱 + 全文检索', ACC_BLUE),
        ('自进化能力', 'GEPA 技能蒸馏，自动沉淀经验', ACC_CYAN),
        ('跨端协作', '桌面 / Web / IM 机器人多端互通', ACC_GREEN),
    ]
    vx = Inches(0.6); vy = Inches(3.7); vw = Inches(2.9); vh = Inches(2.4); vgap = Inches(0.16)
    for i, (t, d, c) in enumerate(visions):
        x = vx + i * (vw + vgap)
        add_round(s, x, vy, vw, vh, fill=BG_CARD, line=LINE_COLOR, radius=0.06)
        add_rect(s, x, vy, vw, Inches(0.08), fill=c)
        add_text(s, x, vy + Inches(0.45), vw, Inches(0.4), t, size=16, bold=True,
                 color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), vy + Inches(1.05), vw - Inches(0.5), Inches(1.2),
                 d, size=11.5, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER, spacing=1.4)
    add_text(s, Inches(0.8), Inches(6.5), Inches(11.73), Inches(0.5),
             '让 AI 从「助手」进化为「同事」与「公司」', size=14, color=ACC_CYAN,
             align=PP_ALIGN.CENTER, italic=True)

    # ---------------- Slide 15: Thanks ----------------
    s = add_slide(prs); set_bg(s, BG_PANEL)
    bg = add_rect(s, 0, 0, PAGE_W, PAGE_H, fill=RGBColor(0x10, 0x10, 0x1A))
    set_gradient(bg, RGBColor(0x16, 0x12, 0x2A), RGBColor(0x08, 0x0A, 0x14), angle=60)
    try:
        s.shapes.add_picture(LOGO, Inches(5.97), Inches(1.5), width=Inches(1.4))
    except Exception as e:
        print('logo skip', e)
    add_text(s, Inches(0.8), Inches(3.15), Inches(11.73), Inches(1.0),
             '感谢关注 Claude Space', size=40, bold=True, color=TEXT_PRIMARY,
             align=PP_ALIGN.CENTER)
    line = add_rect(s, Inches(5.67), Inches(4.25), Inches(2.0), Pt(3), fill=ACC_PURPLE)
    set_gradient(line, ACC_PURPLE, ACC_CYAN, angle=0)
    add_text(s, Inches(0.8), Inches(4.5), Inches(11.73), Inches(0.5),
             'AI 驱动的项目开发工作台', size=20, color=TEXT_SECONDARY,
             align=PP_ALIGN.CENTER)
    # info chips
    infos = ['MIT License', 'Electron + React', '开源项目', 'v1.1.7']
    cx = Inches(3.7); cy = Inches(5.5)
    for info in infos:
        w = Inches(0.2 * len(info) + 0.9)
        add_round(s, cx, cy, w, Inches(0.45), fill=BG_CARD, line=LINE_COLOR, radius=0.5)
        add_text(s, cx, cy, w, Inches(0.45), info, size=11.5, color=TEXT_SECONDARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
        cx = cx + w + Inches(0.2)
    add_text(s, Inches(0.8), Inches(6.5), Inches(11.73), Inches(0.4),
             '多模型  ·  多智能体  ·  多工作空间', size=13, color=TEXT_MUTED,
             align=PP_ALIGN.CENTER, spacing=1.5)

    out = r'E:\claudespace\claude-space\docs\Claude-Space-产品介绍.pptx'
    prs.save(out)
    print('SAVED:', out, 'slides=', len(prs.slides._sldIdLst))
    return out

if __name__ == '__main__':
    build()
